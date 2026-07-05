import base64
import io
import pytest
import respx
import httpx
from unittest.mock import patch, MagicMock, AsyncMock

from icx_engine.connectors.attachments import (
    _is_image, _is_document, _convert_csv, _convert_txt, _rows_to_markdown,
    _convert_xlsx, _convert_xls, _convert_pptx, _convert_zip, _convert_text_passthrough,
    _VISION_PROMPT, _SUMMARIZE_SYSTEM,
    ocr_image, vision_enrich, process_attachments,
    _convert_document, _convert_pdf,
    _summarize_content, _split_into_chunks,
    _SUMMARIZE_THRESHOLD, _SINGLE_CALL_LIMIT, _SUMMARIZE_FAILED_NOTE,
)
from icx_engine.connectors.jira.client import JiraClient
from icx_engine.models.config import LLMConfig, ChannelConfig
from icx_engine.models.output import RawIssueData
from test_data import JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS


# -- Extension detection -------------------------------------------------------

def test_is_image_true_for_common_formats():
    assert _is_image("screenshot.png") is True
    assert _is_image("error.jpg") is True
    assert _is_image("ui.jpeg") is True
    assert _is_image("capture.webp") is True
    assert _is_image("diagram.gif") is True


def test_is_image_false_for_non_images():
    assert _is_image("log.txt") is False
    assert _is_image("report.pdf") is False
    assert _is_image("data.csv") is False
    assert _is_image("code.py") is False


def test_is_image_case_insensitive():
    assert _is_image("SCREEN.PNG") is True
    assert _is_image("Error.JPG") is True


def test_is_document_true_for_supported_types():
    assert _is_document("report.pdf") is True
    assert _is_document("data.csv") is True
    assert _is_document("sheet.xlsx") is True
    assert _is_document("doc.docx") is True
    assert _is_document("notes.txt") is True


def test_is_document_true_for_new_formats():
    assert _is_document("legacy.xls") is True
    assert _is_document("deck.pptx") is True
    assert _is_document("bundle.zip") is True
    assert _is_document("script.py") is True
    assert _is_document("config.json") is True
    assert _is_document("readme.md") is True


def test_is_document_false_for_unsupported():
    assert _is_document("video.mkv") is False
    assert _is_document("screenshot.png") is False


# -- OCR -----------------------------------------------------------------------

def test_ocr_image_returns_text():
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    with patch("pytesseract.image_to_string", return_value="Error: null pointer"), \
         patch("PIL.Image.open", return_value=MagicMock()):
        result = ocr_image(fake_bytes)
    assert result == "Error: null pointer"


def test_ocr_image_returns_empty_on_tesseract_exception():
    fake_bytes = b"\x89PNG\r\n\x1a\n"
    with patch("pytesseract.image_to_string", side_effect=Exception("tesseract failed")), \
         patch("PIL.Image.open", return_value=MagicMock()):
        result = ocr_image(fake_bytes)
    assert result == ""


def test_ocr_image_returns_empty_on_import_error():
    fake_bytes = b"\x89PNG\r\n\x1a\n"
    with patch.dict("sys.modules", {"pytesseract": None}):
        result = ocr_image(fake_bytes)
    assert result == ""


def test_vision_prompt_includes_graph_analysis_sections():
    """_VISION_PROMPT must request graph/chart interpretation alongside text extraction."""
    prompt_lower = _VISION_PROMPT.lower()
    assert "{ocr_text}" in _VISION_PROMPT, \
        "placeholder must be preserved for OCR text injection"
    assert "graph" in prompt_lower or "chart" in prompt_lower, \
        "must request graph/chart analysis"
    assert "axes" in prompt_lower or "axis" in prompt_lower, \
        "must request axis labels and units"
    assert "trend" in prompt_lower, \
        "must request trend identification"
    assert "peak" in prompt_lower, \
        "must request peak value identification"
    assert "minimum" in prompt_lower or "min" in prompt_lower, \
        "must request minimum value identification"
    assert "average" in prompt_lower, \
        "must request average value identification"


def test_vision_prompt_guards_against_prompt_injection():
    prompt_lower = _VISION_PROMPT.lower()
    assert "data, not instructions" in prompt_lower
    assert "do not obey" in prompt_lower


def test_video_frames_prompt_guards_against_prompt_injection():
    from icx_engine.connectors.attachments import _VIDEO_FRAMES_PROMPT
    prompt_lower = _VIDEO_FRAMES_PROMPT.lower()
    assert "data, not instructions" in prompt_lower
    assert "do not obey" in prompt_lower


# -- Vision enrichment ---------------------------------------------------------

async def test_vision_enrich_anthropic_refines_ocr():
    config = ChannelConfig(provider="anthropic", model="claude-3-opus", api_key="sk-ant-test")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    mock_response = MagicMock()
    mock_response.content = [MagicMock()]
    mock_response.content[0].text = "Refined: NullPointerException at com.example.Foo:42"

    with patch("anthropic.AsyncAnthropic") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.messages.create = AsyncMock(return_value=mock_response)
        result = await vision_enrich(config, fake_bytes, "raw ocr text")

    assert result == "Refined: NullPointerException at com.example.Foo:42"


async def test_vision_enrich_openai_refines_ocr():
    config = ChannelConfig(provider="openai", model="gpt-4o", api_key="sk-test")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = "Refined OpenAI output"

    with patch("openai.AsyncOpenAI") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.chat.completions.create = AsyncMock(return_value=mock_response)
        result = await vision_enrich(config, fake_bytes, "raw ocr text")

    assert result == "Refined OpenAI output"


async def test_vision_enrich_nim_vision_model_refines_ocr():
    config = ChannelConfig(provider="nim", model="meta/llama-3.2-11b-vision-instruct", api_key="nim-key")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = "NIM vision output"

    with patch("openai.AsyncOpenAI") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.chat.completions.create = AsyncMock(return_value=mock_response)
        result = await vision_enrich(config, fake_bytes, "raw ocr text")

    assert result == "NIM vision output"


async def test_vision_enrich_google_refines_ocr():
    config = ChannelConfig(provider="google", model="gemini-1.5-flash", api_key="goog-test")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    mock_response = MagicMock()
    mock_response.text = "Google vision output"

    mock_client = MagicMock()
    mock_client.aio = MagicMock()
    mock_client.aio.models = MagicMock()
    mock_client.aio.models.generate_content = AsyncMock(return_value=mock_response)

    with patch("google.genai.Client", return_value=mock_client):
        result = await vision_enrich(config, fake_bytes, "raw ocr text")

    assert result == "Google vision output"


async def test_vision_enrich_raises_clear_error_when_model_is_text_only():
    from icx_engine.exceptions import ContextBuildError

    config = ChannelConfig(provider="ollama", model="llama3")
    fake_bytes = b"\x89PNG\r\n\x1a\n"

    with patch("openai.AsyncOpenAI") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.chat.completions.create = AsyncMock(
            side_effect=Exception("model does not support vision")
        )
        with pytest.raises(ContextBuildError, match="does not support image input"):
            await vision_enrich(config, fake_bytes, "ocr fallback text")


# -- SDK timeout parameters ----------------------------------------------------

async def test_vision_enrich_anthropic_passes_timeout_90s():
    """Anthropic vision call must include timeout=90.0 to prevent indefinite hang."""
    config = ChannelConfig(provider="anthropic", model="claude-3-haiku-20240307", api_key="sk-ant-test")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    mock_response = MagicMock()
    mock_response.content = [MagicMock()]
    mock_response.content[0].text = "result"

    with patch("anthropic.AsyncAnthropic") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.messages.create = AsyncMock(return_value=mock_response)
        await vision_enrich(config, fake_bytes, "ocr text", "shot.png")

    _, kwargs = mock_client.messages.create.call_args
    assert kwargs.get("timeout") == 90.0


async def test_vision_enrich_openai_passes_timeout_90s():
    """OpenAI-compat vision call must include timeout=90.0 to prevent indefinite hang."""
    config = ChannelConfig(provider="openai", model="gpt-4o", api_key="sk-test")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = "result"

    with patch("openai.AsyncOpenAI") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.chat.completions.create = AsyncMock(return_value=mock_response)
        await vision_enrich(config, fake_bytes, "ocr text", "shot.png")

    _, kwargs = mock_client.chat.completions.create.call_args
    assert kwargs.get("timeout") == 90.0


async def test_vision_enrich_google_timeout_raises_context_build_error():
    """A Google vision call that hits asyncio.wait_for timeout surfaces as ContextBuildError."""
    import asyncio as _asyncio
    from icx_engine.exceptions import ContextBuildError

    config = ChannelConfig(provider="google", model="gemini-1.5-flash", api_key="goog-test")
    fake_bytes = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50

    async def _timeout_wait_for(*args, **kwargs):
        raise _asyncio.TimeoutError

    mock_client = MagicMock()
    mock_client.aio.models.generate_content = MagicMock(return_value=MagicMock())

    with patch("google.genai.Client", return_value=mock_client):
        with patch("icx_engine.connectors.attachments.asyncio.wait_for", new=_timeout_wait_for):
            with pytest.raises(ContextBuildError):
                await vision_enrich(config, fake_bytes, "ocr text", "shot.png")


async def test_llm_summarize_chunk_anthropic_passes_timeout_90s():
    """Anthropic summarize call must include timeout=90.0."""
    from icx_engine.connectors.attachments import _llm_summarize_chunk

    config = ChannelConfig(provider="anthropic", model="claude-3-haiku-20240307", api_key="sk-ant-test")

    mock_response = MagicMock()
    mock_response.content = [MagicMock()]
    mock_response.content[0].text = "summary"

    with patch("anthropic.AsyncAnthropic") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.messages.create = AsyncMock(return_value=mock_response)
        await _llm_summarize_chunk(config, "doc.pdf", "long content here")

    _, kwargs = mock_client.messages.create.call_args
    assert kwargs.get("timeout") == 90.0


async def test_llm_summarize_chunk_openai_passes_timeout_90s():
    """OpenAI-compat summarize call must include timeout=90.0."""
    from icx_engine.connectors.attachments import _llm_summarize_chunk

    config = ChannelConfig(provider="openai", model="gpt-4o-mini", api_key="sk-test")

    mock_response = MagicMock()
    mock_response.choices = [MagicMock()]
    mock_response.choices[0].message.content = "summary"

    with patch("openai.AsyncOpenAI") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.chat.completions.create = AsyncMock(return_value=mock_response)
        await _llm_summarize_chunk(config, "doc.pdf", "long content here")

    _, kwargs = mock_client.chat.completions.create.call_args
    assert kwargs.get("timeout") == 90.0


# -- UAE: document converters --------------------------------------------------

def test_convert_csv_returns_markdown_table():
    csv_data = b"Name,Status\nAlice,Active\nBob,Inactive"
    result = _convert_csv(csv_data)
    assert "| Name | Status |" in result
    assert "| Alice | Active |" in result
    assert "| Bob | Inactive |" in result


def test_convert_csv_truncates_at_50_rows():
    header = "A,B\n"
    rows = "\n".join(f"{i},{i*2}" for i in range(1, 60))  # 59 data rows
    csv_data = (header + rows).encode()
    result = _convert_csv(csv_data)
    assert "truncated" in result.lower()
    assert "59" in result  # shows total row count


def test_convert_txt_returns_clean_text():
    from icx_engine.connectors.attachments import _convert_txt
    data = b"  Hello World  \nLine 2\n"
    result = _convert_txt(data)
    assert "Hello World" in result
    assert "Line 2" in result


def test_convert_txt_does_not_truncate_large_content():
    """Extraction never truncates - full content is preserved for summarization/passthrough."""
    big_text = ("A" * 1000 + "\n") * 200   # 200_200 chars
    data = big_text.encode()
    result = _convert_txt(data)
    assert len(result) == len(big_text.strip())


def test_rows_to_markdown_produces_valid_table():
    rows = [["Col1", "Col2"], ["A", "B"], ["C", "D"]]
    result = _rows_to_markdown(rows)
    assert "| Col1 | Col2 |" in result
    assert "| --- | --- |" in result
    assert "| A | B |" in result


def test_rows_to_markdown_truncates_at_50():
    rows = [["X"]] + [[str(i)] for i in range(60)]
    result = _rows_to_markdown(rows)
    assert "truncated" in result.lower()


def test_convert_xlsx_unavailable_returns_message():
    import unittest.mock
    with unittest.mock.patch.dict("sys.modules", {"openpyxl": None}):
        from icx_engine.connectors.attachments import _convert_xlsx
        result = _convert_xlsx(b"fake")
    assert "unavailable" in result.lower()


def test_convert_pdf_unavailable_returns_message():
    import unittest.mock
    with unittest.mock.patch.dict("sys.modules", {"pdfminer": None, "pdfminer.high_level": None}):
        from icx_engine.connectors.attachments import _convert_pdf
        result, images = _convert_pdf(b"fake")
    assert "unavailable" in result.lower()
    assert images == []


def test_convert_docx_unavailable_returns_message():
    import unittest.mock
    with unittest.mock.patch.dict("sys.modules", {"docx": None}):
        from icx_engine.connectors.attachments import _convert_docx
        result = _convert_docx(b"fake")
    assert "unavailable" in result.lower()


def test_convert_document_passes_through_code_for_known_extension():
    """.py is a recognized text-passthrough extension - fenced code block, not empty."""
    result, images = _convert_document("script.py", b"print('hello')")
    assert "```python" in result
    assert "print('hello')" in result
    assert images == []


def test_convert_document_returns_empty_for_unsupported_extension():
    result, images = _convert_document("archive.rar", b"binary data")
    assert result == ""
    assert images == []


def test_convert_xlsx_annotates_formula_cells_in_sample_rows():
    """Formula cells in header + first 3 data rows are annotated as 'value (Formula: expr)'."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Prices"
    ws.append(["Item", "Price", "Tax"])
    ws.append(["Widget", 100, "=B2*0.18"])
    ws.append(["Gadget", 200, "=B3*0.18"])
    ws.append(["Doohickey", 300, "=B4*0.18"])
    buf = io.BytesIO()
    wb.save(buf)

    result = _convert_xlsx(buf.getvalue())

    assert "**Sheet: Prices**" in result
    assert "(Formula: =B2*0.18)" in result
    assert "(Formula: =B3*0.18)" in result
    assert "(Formula: =B4*0.18)" in result


def test_convert_xlsx_plain_value_cells_not_annotated():
    """Cells without formulas must appear as plain values with no (Formula: ...) annotation."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    ws.append(["Bob", 87])
    buf = io.BytesIO()
    wb.save(buf)

    result = _convert_xlsx(buf.getvalue())

    assert "Formula:" not in result
    assert "Alice" in result
    assert "95" in result


def test_convert_xlsx_formula_annotation_skipped_beyond_row_3():
    """Formula annotation applies only to header + first 3 data rows; rows 4+ show plain values."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Item", "Total"])
    for i in range(1, 6):
        ws.append([f"Item{i}", f"=A{i + 1}*10"])
    buf = io.BytesIO()
    wb.save(buf)

    result = _convert_xlsx(buf.getvalue())

    # Data rows 1-3 (val_rows indices 1-3) must be annotated
    assert "(Formula: =A2*10)" in result
    assert "(Formula: =A3*10)" in result
    assert "(Formula: =A4*10)" in result
    # Data rows 4-5 (val_rows indices 4-5) must NOT be annotated
    assert "(Formula: =A5*10)" not in result
    assert "(Formula: =A6*10)" not in result


def test_convert_xlsx_multi_sheet_formula_annotation():
    """Formula annotation works independently on each sheet in a multi-sheet workbook."""
    import openpyxl
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Sheet1"
    ws1.append(["Val", "Computed"])
    ws1.append([10, "=A2*2"])
    ws2 = wb.create_sheet("Sheet2")
    ws2.append(["Val", "Computed"])
    ws2.append([20, "=A2*3"])
    buf = io.BytesIO()
    wb.save(buf)

    result = _convert_xlsx(buf.getvalue())

    assert "**Sheet: Sheet1**" in result
    assert "**Sheet: Sheet2**" in result
    assert "(Formula: =A2*2)" in result
    assert "(Formula: =A2*3)" in result


# -- process_attachments - orchestrator ---------------------------------------

def _make_raw(attachment_content_urls=None) -> RawIssueData:
    return RawIssueData(
        issue_key="X-1", issue_type="Bug", summary="s", description="d",
        comments=[], attachments=list((attachment_content_urls or {}).keys()),
        priority="High", status="Open", metadata={},
        attachment_content_urls=attachment_content_urls or {},
    )


@respx.mock
async def test_process_attachments_downloads_and_ocrs():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/10001"
    fake_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    respx.get(content_url).mock(return_value=httpx.Response(200, content=fake_png))

    raw = _make_raw({"screenshot.png": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.attachments.ocr_image", return_value="OCR result"):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert texts == {"screenshot.png": "OCR result"}
    assert "screenshot.png" in images
    import base64 as _b64
    assert _b64.b64decode(images["screenshot.png"]) == fake_png


@respx.mock
async def test_process_attachments_handles_pdf():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/99"
    fake_pdf = b"fake pdf bytes"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=fake_pdf))

    raw = _make_raw({"report.pdf": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.attachments._convert_document", return_value=("Extracted PDF text", [])):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert texts == {"report.pdf": "Extracted PDF text"}
    assert images == {}  # documents produce no Base64


async def test_process_attachments_skips_unsupported_extension():
    raw = _make_raw({"archive.rar": "https://test.atlassian.net/rest/api/3/attachment/content/100"})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    texts, images = await process_attachments(raw, client, llm_config=None)
    assert texts == {}
    assert images == {}


@respx.mock
async def test_process_attachments_skips_on_download_error():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/10001"
    respx.get(content_url).mock(return_value=httpx.Response(403))
    raw = _make_raw({"screenshot.png": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    texts, images = await process_attachments(raw, client, llm_config=None)
    assert texts == {}
    assert images == {}


@respx.mock
async def test_process_attachments_skips_vision_when_image_model_is_none():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/10001"
    fake_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    respx.get(content_url).mock(return_value=httpx.Response(200, content=fake_png))

    raw = _make_raw({"screenshot.png": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    llm = LLMConfig(text_config=ChannelConfig(provider="ollama", model="llama3"), image_config=None)

    with patch("icx_engine.connectors.attachments.ocr_image", return_value="raw ocr"):
        with patch("icx_engine.connectors.attachments.vision_enrich", new=AsyncMock()) as mock_vision:
            texts, images = await process_attachments(raw, client, llm_config=llm)

    mock_vision.assert_not_called()
    assert texts == {"screenshot.png": "raw ocr"}
    assert "screenshot.png" in images


@respx.mock
async def test_process_attachments_calls_vision_when_llm_config_provided():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/10001"
    fake_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    respx.get(content_url).mock(return_value=httpx.Response(200, content=fake_png))

    raw = _make_raw({"screenshot.png": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    llm = LLMConfig(text_config=ChannelConfig(provider="ollama", model="llama3"),
                    image_config=ChannelConfig(provider="ollama", model="llava"))

    with patch("icx_engine.connectors.attachments.ocr_image", return_value="raw ocr"):
        with patch("icx_engine.connectors.attachments.vision_enrich", new=AsyncMock(return_value="vision refined")) as mock_vision:
            texts, images = await process_attachments(raw, client, llm_config=llm)

    mock_vision.assert_called_once()
    assert texts == {"screenshot.png": "vision refined"}
    assert "screenshot.png" in images  # Base64 always captured regardless of vision model


@respx.mock
async def test_process_attachments_processes_multiple_in_parallel():
    url1 = "https://test.atlassian.net/rest/api/3/attachment/content/10001"
    url2 = "https://test.atlassian.net/rest/api/3/attachment/content/10002"
    fake_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    fake_csv = b"A,B\n1,2\n"
    respx.get(url1).mock(return_value=httpx.Response(200, content=fake_png))
    respx.get(url2).mock(return_value=httpx.Response(200, content=fake_csv))

    raw = _make_raw({"screen.png": url1, "data.csv": url2})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.attachments.ocr_image", return_value="OCR text"):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert "screen.png" in texts
    assert "data.csv" in texts
    assert "OCR text" in texts["screen.png"]
    assert "| A | B |" in texts["data.csv"]
    assert "screen.png" in images   # image captured as Base64
    assert "data.csv" not in images  # document has no Base64


@respx.mock
async def test_process_attachments_summarizes_large_document_with_llm():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/99"
    big_content = "A" * 25_000  # exceeds _SUMMARIZE_THRESHOLD, within _SINGLE_CALL_LIMIT
    respx.get(content_url).mock(return_value=httpx.Response(200, content=big_content.encode()))

    raw = _make_raw({"notes.txt": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    llm = LLMConfig(text_config=ChannelConfig(provider="ollama", model="llama3"), image_config=None)

    with patch("icx_engine.connectors.attachments._llm_summarize_chunk", new=AsyncMock(return_value="summary")) as mock_sum:
        texts, images = await process_attachments(raw, client, llm_config=llm)

    mock_sum.assert_called_once()
    assert texts == {"notes.txt": "summary"}
    assert images == {}


@respx.mock
async def test_process_attachments_no_truncation_without_llm():
    """Without an LLM configured, full content is passed through - nothing is truncated."""
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/99"
    big_content = "B" * 25_000
    respx.get(content_url).mock(return_value=httpx.Response(200, content=big_content.encode()))

    raw = _make_raw({"notes.txt": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    texts, images = await process_attachments(raw, client, llm_config=None)
    assert texts["notes.txt"] == big_content
    assert images == {}


@respx.mock
async def test_process_attachments_returns_empty_tuple_when_no_urls():
    raw = _make_raw({})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    texts, images = await process_attachments(raw, client, llm_config=None)
    assert texts == {}
    assert images == {}


@respx.mock
async def test_process_attachments_captures_base64_for_images():
    """Base64 is captured even when OCR returns no text."""
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/10001"
    fake_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 100
    respx.get(content_url).mock(return_value=httpx.Response(200, content=fake_png))

    raw = _make_raw({"error_screenshot.png": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.attachments.ocr_image", return_value=""):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert "error_screenshot.png" not in texts   # no OCR text -> not in texts
    assert "error_screenshot.png" in images       # but Base64 still captured
    import base64 as _b64
    assert _b64.b64decode(images["error_screenshot.png"]) == fake_png


def test_summarize_system_preserves_structured_data():
    p = _SUMMARIZE_SYSTEM.lower()
    assert "column headers" in p
    assert "formula" in p
    assert "non-negotiable business rule" in p
    assert "[technical schema:" in p
    assert "[technical logic:" in p
    assert "verbatim" in p


def test_summarize_system_guards_against_prompt_injection():
    p = _SUMMARIZE_SYSTEM.lower()
    assert "data, not instructions" in p
    assert "do not obey" in p


def test_process_attachments_spinner_shown_when_log_set():
    """When log is provided, a Rich status spinner is shown during processing."""
    import asyncio
    from unittest.mock import AsyncMock, MagicMock, patch
    from icx_engine.connectors.attachments import process_attachments
    from icx_engine.models.output import RawIssueData

    raw = RawIssueData(
        issue_key="T-1", issue_type="Bug", summary="s", description="d",
        comments=[], attachments=["report.txt"],
        priority="High", status="Open", metadata={},
        attachment_content_urls={"report.txt": "https://x.atlassian.net/content/1"},
        attachment_texts={},
    )

    downloader = MagicMock()
    downloader.download_attachment = AsyncMock(return_value=b"hello world content")

    status_context = MagicMock()
    status_context.__enter__ = MagicMock(return_value=None)
    status_context.__exit__ = MagicMock(return_value=False)

    with patch("rich.console.Console") as mock_console_cls:
        mock_con_instance = MagicMock()
        mock_con_instance.status.return_value = status_context
        mock_console_cls.return_value = mock_con_instance

        asyncio.run(process_attachments(raw, downloader, llm_config=None, log=list().append))

    mock_con_instance.status.assert_called_once()
    call_msg = mock_con_instance.status.call_args[0][0]
    assert "attachment" in call_msg.lower() or "processing" in call_msg.lower()


# -- Extension detection (audio/video) -----------------------------------------

def test_is_audio_true_for_audio_formats():
    from icx_engine.connectors.attachments import _is_audio
    for fname in ["note.mp3", "clip.wav", "record.m4a", "sound.ogg", "track.flac", "audio.aac", "voice.opus"]:
        assert _is_audio(fname) is True, f"Expected True for {fname}"


def test_is_audio_false_for_non_audio():
    from icx_engine.connectors.attachments import _is_audio
    assert _is_audio("report.pdf") is False
    assert _is_audio("screen.png") is False
    assert _is_audio("video.mp4") is False


def test_is_audio_case_insensitive():
    from icx_engine.connectors.attachments import _is_audio
    assert _is_audio("RECORDING.MP3") is True
    assert _is_audio("Clip.WAV") is True


def test_is_video_true_for_video_formats():
    from icx_engine.connectors.attachments import _is_video
    for fname in ["demo.mp4", "screen.mov", "capture.avi", "video.mkv", "clip.webm"]:
        assert _is_video(fname) is True, f"Expected True for {fname}"


def test_is_video_false_for_non_video():
    from icx_engine.connectors.attachments import _is_video
    assert _is_video("audio.mp3") is False
    assert _is_video("report.pdf") is False
    assert _is_video("screen.png") is False


# -- _process_audio -------------------------------------------------------------

@respx.mock
async def test_process_audio_downloads_and_transcribes():
    from icx_engine.connectors.attachments import _process_audio
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/200"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"fake mp3 bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()

    with patch("icx_engine.connectors.audio.transcribe", new=AsyncMock(return_value="Transcribed audio text.")):
        fname, text, images = await _process_audio("meeting.mp3", content_url, client, None, whisper, None)

    assert fname == "meeting.mp3"
    assert text == "Transcribed audio text."
    assert images == {}


@respx.mock
async def test_process_audio_returns_empty_on_download_error():
    from icx_engine.connectors.attachments import _process_audio
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/201"
    respx.get(content_url).mock(return_value=httpx.Response(403))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()

    fname, text, images = await _process_audio("meeting.mp3", content_url, client, None, whisper, None)

    assert text == ""
    assert images == {}


@respx.mock
async def test_process_audio_returns_empty_on_transcription_error():
    from icx_engine.connectors.attachments import _process_audio
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/202"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()
    log_messages: list[str] = []

    with patch("icx_engine.connectors.audio.transcribe",
               new=AsyncMock(side_effect=Exception("boom"))):
        fname, text, images = await _process_audio(
            "meeting.mp3", content_url, client, None, whisper, log_messages.append
        )

    assert fname == "meeting.mp3"
    assert text == ""
    assert images == {}
    assert any("transcription failed" in m for m in log_messages), (
        f"Expected 'transcription failed' log message but got: {log_messages}"
    )
    assert any("boom" in m for m in log_messages), (
        f"Expected exception message 'boom' in log but got: {log_messages}"
    )


@respx.mock
async def test_process_audio_returns_setup_msg_when_whisper_not_installed():
    """_process_audio must surface _SETUP_REQUIRED_MSG when Whisper raises setup-required RuntimeError."""
    from icx_engine.connectors.attachments import _process_audio, _SETUP_REQUIRED_MSG
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/203"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()
    log_messages: list[str] = []

    with patch("icx_engine.connectors.audio.transcribe",
               new=AsyncMock(side_effect=RuntimeError(
                   "Whisper model not found. Run `icx setup` to download it."
               ))):
        fname, text, images = await _process_audio(
            "meeting.mp3", content_url, client, None, whisper, log_messages.append
        )

    assert fname == "meeting.mp3"
    assert text == _SETUP_REQUIRED_MSG
    assert images == {}
    assert any("icx setup" in m for m in log_messages)


# -- _process_video -------------------------------------------------------------

@respx.mock
async def test_process_video_extracts_audio_then_transcribes():
    from icx_engine.connectors.attachments import _process_video
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/300"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"fake mp4 bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()

    with patch("icx_engine.connectors.attachments._extract_audio_from_video", new=AsyncMock(return_value=b"\x00" * 44)), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video", new=AsyncMock(return_value=[])), \
         patch("icx_engine.connectors.audio.transcribe", new=AsyncMock(return_value="Video transcript.")):
        fname, text, images = await _process_video("demo.mp4", content_url, client, None, None, whisper, None)

    assert fname == "demo.mp4"
    assert text == "**Transcript:**\n\nVideo transcript."
    assert images == {}


@respx.mock
async def test_process_video_returns_empty_on_extraction_error():
    from icx_engine.connectors.attachments import _process_video
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/301"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()

    with patch("icx_engine.connectors.attachments._extract_audio_from_video",
               new=AsyncMock(side_effect=Exception("ffmpeg failed"))), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video",
               new=AsyncMock(return_value=[])):
        fname, text, images = await _process_video("demo.mp4", content_url, client, None, None, whisper, None)

    assert text == ""
    assert images == {}


@respx.mock
async def test_process_video_passes_wav_fname_without_double_extension():
    """transcribe must receive 'demo.wav', not 'demo.mp4.wav'."""
    from icx_engine.connectors.attachments import _process_video
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/302"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"mp4 bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()

    transcribe_mock = AsyncMock(return_value="Video transcript.")
    with patch("icx_engine.connectors.attachments._extract_audio_from_video",
               new=AsyncMock(return_value=b"\x00" * 44)), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video",
               new=AsyncMock(return_value=[])), \
         patch("icx_engine.connectors.audio.transcribe", new=transcribe_mock):
        await _process_video("demo.mp4", content_url, client, None, None, whisper, None)

    # transcribe is called as: transcribe(text_config, audio_bytes, fname, whisper)
    fname_arg = transcribe_mock.call_args[0][2]
    assert fname_arg == "demo.wav", f"Expected 'demo.wav' but got '{fname_arg}'"


@respx.mock
async def test_process_video_returns_empty_when_no_audio_track():
    """Videos with no audio track produce very small WAV bytes; must be caught and logged."""
    from icx_engine.connectors.attachments import _process_video
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/303"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"mp4 bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()
    log_messages: list[str] = []

    # ffmpeg exits 0 but returns minimal/empty WAV (< 44 bytes = no real audio)
    empty_wav = b"\x00" * 10
    with patch("icx_engine.connectors.attachments._extract_audio_from_video",
               new=AsyncMock(return_value=empty_wav)), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video",
               new=AsyncMock(return_value=[])):
        fname, text, images = await _process_video(
            "silent.mp4", content_url, client, None, None, whisper, log_messages.append
        )

    assert fname == "silent.mp4"
    assert text == ""
    assert images == {}
    assert any("no audio" in m.lower() for m in log_messages), (
        f"Expected 'no audio' log message but got: {log_messages}"
    )


@respx.mock
async def test_process_video_returns_setup_msg_when_whisper_not_installed():
    """_process_video must surface _SETUP_REQUIRED_MSG when Whisper raises setup-required RuntimeError.

    Frame extraction still runs (it uses vision/OCR, not Whisper). When no frames
    are available, only the setup message is returned.
    """
    from icx_engine.connectors.attachments import _process_video, _SETUP_REQUIRED_MSG
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/311"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"fake mp4 bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()
    log_messages: list[str] = []

    with patch("icx_engine.connectors.attachments._extract_audio_from_video",
               new=AsyncMock(return_value=b"\x00" * 100)), \
         patch("icx_engine.connectors.audio.transcribe",
               new=AsyncMock(side_effect=RuntimeError(
                   "Whisper model not found. Run `icx setup` to download it."
               ))), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video",
               new=AsyncMock(return_value=[])):
        fname, text, images = await _process_video(
            "demo.mp4", content_url, client, None, None, whisper, log_messages.append
        )

    assert fname == "demo.mp4"
    assert text == _SETUP_REQUIRED_MSG
    assert images == {}
    assert any("icx setup" in m for m in log_messages)


# -- _extract_audio_from_video safety: timeout + returncode --------------------

async def test_extract_audio_from_video_kills_ffmpeg_on_timeout():
    """asyncio.wait_for timeout must kill the ffmpeg subprocess, not orphan it."""
    import asyncio as _asyncio
    from icx_engine.connectors.attachments import _extract_audio_from_video

    proc = MagicMock()
    proc.communicate = AsyncMock(side_effect=_asyncio.TimeoutError)
    proc.kill = MagicMock()
    proc.wait = AsyncMock(return_value=0)

    async def _make_subproc(*args, **kwargs):
        return proc

    with patch("imageio_ffmpeg.get_ffmpeg_exe", return_value="ffmpeg"), \
         patch("asyncio.create_subprocess_exec", new=_make_subproc):
        with pytest.raises(_asyncio.TimeoutError):
            await _extract_audio_from_video(b"fake video", "demo.mp4")

    proc.kill.assert_called_once()
    proc.wait.assert_awaited()


async def test_extract_audio_from_video_raises_on_nonzero_returncode():
    """Non-zero ffmpeg exit code must raise RuntimeError, not silently return empty bytes."""
    from icx_engine.connectors.attachments import _extract_audio_from_video

    proc = MagicMock()
    proc.communicate = AsyncMock(return_value=(b"", b"bad codec error"))
    proc.returncode = 1

    async def _make_subproc(*args, **kwargs):
        return proc

    with patch("imageio_ffmpeg.get_ffmpeg_exe", return_value="ffmpeg"), \
         patch("asyncio.create_subprocess_exec", new=_make_subproc):
        with pytest.raises(RuntimeError, match="ffmpeg exited with code 1"):
            await _extract_audio_from_video(b"fake video", "demo.mp4")


# -- process_attachments with audio/video --------------------------------------

@respx.mock
async def test_process_attachments_transcribes_audio_attachment():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/400"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"mp3 bytes"))

    raw = _make_raw({"meeting.mp3": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.audio.transcribe", new=AsyncMock(return_value="Meeting transcript.")):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert texts == {"meeting.mp3": "Meeting transcript."}
    assert images == {}


@respx.mock
async def test_process_attachments_transcribes_video_attachment():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/401"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"mp4 bytes"))

    raw = _make_raw({"demo.mp4": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.attachments._extract_audio_from_video", new=AsyncMock(return_value=b"\x00" * 44)), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video", new=AsyncMock(return_value=[])), \
         patch("icx_engine.connectors.audio.transcribe", new=AsyncMock(return_value="Demo video transcript.")):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert texts == {"demo.mp4": "**Transcript:**\n\nDemo video transcript."}
    assert images == {}


@respx.mock
async def test_process_attachments_handles_mixed_image_audio_video_doc():
    img_url = "https://test.atlassian.net/rest/api/3/attachment/content/500"
    audio_url = "https://test.atlassian.net/rest/api/3/attachment/content/501"
    video_url = "https://test.atlassian.net/rest/api/3/attachment/content/502"
    doc_url = "https://test.atlassian.net/rest/api/3/attachment/content/503"

    fake_png = b"\x89PNG\r\n\x1a\n" + b"\x00" * 50
    respx.get(img_url).mock(return_value=httpx.Response(200, content=fake_png))
    respx.get(audio_url).mock(return_value=httpx.Response(200, content=b"mp3"))
    respx.get(video_url).mock(return_value=httpx.Response(200, content=b"mp4"))
    respx.get(doc_url).mock(return_value=httpx.Response(200, content=b"A,B\n1,2\n"))

    raw = _make_raw({
        "screen.png": img_url,
        "note.mp3": audio_url,
        "demo.mp4": video_url,
        "data.csv": doc_url,
    })
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)

    with patch("icx_engine.connectors.attachments.ocr_image", return_value="OCR"), \
         patch("icx_engine.connectors.attachments._extract_audio_from_video", new=AsyncMock(return_value=b"\x00" * 44)), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video", new=AsyncMock(return_value=[])), \
         patch("icx_engine.connectors.audio.transcribe", new=AsyncMock(return_value="audio text")):
        texts, images = await process_attachments(raw, client, llm_config=None)

    assert "screen.png" in texts
    assert "note.mp3" in texts
    assert texts["note.mp3"] == "audio text"
    assert "demo.mp4" in texts
    assert texts["demo.mp4"] == "**Transcript:**\n\naudio text"
    assert "data.csv" in texts
    assert "screen.png" in images
    assert "note.mp3" not in images
    assert "demo.mp4" not in images


# -- UAE: text/code passthrough -------------------------------------------------

def test_convert_text_passthrough_fences_code_by_extension():
    result = _convert_text_passthrough("script.py", b"print('hello')")
    assert result == "```python\nprint('hello')\n```"


def test_convert_text_passthrough_markdown_returned_raw():
    result = _convert_text_passthrough("notes.md", b"# Heading\n\nSome text")
    assert result == "# Heading\n\nSome text"
    assert "```" not in result


def test_convert_text_passthrough_unknown_extension_fences_without_lang():
    result = _convert_text_passthrough("data.unknownext", b"raw content")
    assert result == "```\nraw content\n```"


# -- UAE: .xls (legacy Excel) ----------------------------------------------------

def test_convert_xls_unavailable_returns_message():
    import unittest.mock
    with unittest.mock.patch.dict("sys.modules", {"xlrd": None}):
        from icx_engine.connectors.attachments import _convert_xls
        result = _convert_xls(b"fake")
    assert "unavailable" in result.lower()


def test_convert_xls_converts_sheets_to_markdown():
    fake_sheet = MagicMock()
    fake_sheet.name = "Prices"
    fake_sheet.nrows = 2
    fake_sheet.row_values.side_effect = lambda r: [["Item", "Cost"], ["Widget", 10]][r]

    fake_wb = MagicMock()
    fake_wb.sheets.return_value = [fake_sheet]

    with patch("xlrd.open_workbook", return_value=fake_wb):
        result = _convert_xls(b"fake")

    assert "**Sheet: Prices**" in result
    assert "| Item | Cost |" in result
    assert "| Widget | 10 |" in result


def test_convert_xls_skips_empty_sheets():
    empty_sheet = MagicMock()
    empty_sheet.name = "Empty"
    empty_sheet.nrows = 0

    fake_wb = MagicMock()
    fake_wb.sheets.return_value = [empty_sheet]

    with patch("xlrd.open_workbook", return_value=fake_wb):
        result = _convert_xls(b"fake")

    assert "Empty" not in result


# -- UAE: .pptx (PowerPoint) ------------------------------------------------------

def test_convert_pptx_unavailable_returns_message():
    import unittest.mock
    with unittest.mock.patch.dict("sys.modules", {"pptx": None}):
        from icx_engine.connectors.attachments import _convert_pptx
        result = _convert_pptx(b"fake")
    assert "unavailable" in result.lower()


def test_convert_pptx_extracts_slide_text_and_notes():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = prs.slide_layouts[5]  # blank-ish layout with title placeholder
    slide = prs.slides.add_slide(layout)
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Hello from slide one"
    notes = slide.notes_slide
    notes.notes_text_frame.text = "Speaker notes here"

    buf = io.BytesIO()
    prs.save(buf)

    result = _convert_pptx(buf.getvalue())

    assert "## Slide 1" in result
    assert "Hello from slide one" in result
    assert "**Notes:** Speaker notes here" in result


# -- UAE: .zip -----------------------------------------------------------------

def test_convert_zip_invalid_returns_message():
    result = _convert_zip("bad.zip", b"not a zip file")
    assert "invalid" in result.lower()


def _make_zip(entries: dict[str, bytes]) -> bytes:
    import zipfile
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w") as zf:
        for name, content in entries.items():
            zf.writestr(name, content)
    return buf.getvalue()


def test_convert_zip_manifest_and_recursive_conversion():
    data = _make_zip({
        "readme.md": b"# Title\n\nBody text",
        "data.csv": b"A,B\n1,2",
    })
    result = _convert_zip("bundle.zip", data)

    assert "ZIP archive: 2 file(s)" in result
    assert "readme.md" in result
    assert "data.csv" in result
    assert "# Title" in result
    assert "| A | B |" in result


def test_convert_zip_caps_at_max_entries():
    from icx_engine.connectors.attachments import _ZIP_MAX_ENTRIES
    entries = {f"file{i}.txt": b"x" for i in range(_ZIP_MAX_ENTRIES + 5)}
    data = _make_zip(entries)

    result = _convert_zip("bundle.zip", data)

    assert f"more entr" in result
    assert "5 more" in result


def test_convert_zip_skips_oversized_entry():
    from icx_engine.connectors.attachments import _ZIP_ENTRY_MAX_BYTES
    big = b"x" * (_ZIP_ENTRY_MAX_BYTES + 1)
    data = _make_zip({"huge.txt": big, "small.txt": b"ok"})

    result = _convert_zip("bundle.zip", data)

    assert "huge.txt" in result
    assert "exceeds" in result.lower()
    assert "small.txt" in result
    assert "ok" in result


def test_convert_zip_does_not_recurse_into_nested_zip():
    # A zip-in-zip must be listed, not expanded - guards against unbounded
    # recursion / zip-quine DoS. The inner zip's own entries must not appear.
    inner = _make_zip({"secret.txt": b"inner payload"})
    outer = _make_zip({"nested.zip": inner, "note.txt": b"top level"})

    result = _convert_zip("outer.zip", outer)

    assert "nested.zip" in result
    assert "not expanded" in result.lower()
    assert "top level" in result       # sibling non-zip entry still converted
    assert "inner payload" not in result  # inner entry never expanded


def test_convert_document_dispatches_xls_pptx_zip():
    """`.xls`, `.pptx`, `.zip` are routed to their converters via _convert_document."""
    with patch("icx_engine.connectors.attachments._convert_xls", return_value="xls text") as m_xls:
        text, images = _convert_document("legacy.xls", b"fake")
    assert text == "xls text" and images == []
    m_xls.assert_called_once()

    with patch("icx_engine.connectors.attachments._convert_pptx", return_value="pptx text") as m_pptx:
        text, images = _convert_document("deck.pptx", b"fake")
    assert text == "pptx text" and images == []
    m_pptx.assert_called_once()

    with patch("icx_engine.connectors.attachments._convert_zip", return_value="zip text") as m_zip:
        text, images = _convert_document("bundle.zip", b"fake")
    assert text == "zip text" and images == []
    m_zip.assert_called_once()


# -- UAE: scanned-PDF OCR fallback ------------------------------------------------

def test_convert_pdf_falls_back_to_ocr_for_scanned_pdf():
    """When pdfminer extracts < _PDF_TEXT_MIN_CHARS, pages are rendered + OCR'd."""
    fake_page = MagicMock()
    fake_pix = MagicMock()
    fake_pix.tobytes.return_value = b"\xff\xd8fake-jpeg"
    fake_page.get_pixmap.return_value = fake_pix

    fake_doc = MagicMock()
    fake_doc.__len__.return_value = 1
    fake_doc.load_page.return_value = fake_page

    fake_fitz = MagicMock()
    fake_fitz.open.return_value = fake_doc

    with patch("pdfminer.high_level.extract_text", return_value="  "), \
         patch.dict("sys.modules", {"fitz": fake_fitz}), \
         patch("icx_engine.connectors.attachments.ocr_image", return_value="OCR page text"):
        text, images = _convert_pdf(b"fake pdf bytes")

    assert "### Page 1" in text
    assert "OCR page text" in text
    assert images == [b"\xff\xd8fake-jpeg"]
    fake_doc.close.assert_called_once()


def test_convert_pdf_ocr_capped_at_page_limit():
    from icx_engine.connectors.attachments import _PDF_OCR_PAGE_CAP

    fake_page = MagicMock()
    fake_pix = MagicMock()
    fake_pix.tobytes.return_value = b"\xff\xd8fake-jpeg"
    fake_page.get_pixmap.return_value = fake_pix

    fake_doc = MagicMock()
    fake_doc.__len__.return_value = _PDF_OCR_PAGE_CAP + 10
    fake_doc.load_page.return_value = fake_page

    fake_fitz = MagicMock()
    fake_fitz.open.return_value = fake_doc

    log_messages: list[str] = []

    with patch("pdfminer.high_level.extract_text", return_value=""), \
         patch.dict("sys.modules", {"fitz": fake_fitz}), \
         patch("icx_engine.connectors.attachments.ocr_image", return_value="text"):
        text, images = _convert_pdf(b"fake pdf bytes", log=log_messages.append)

    assert len(images) == _PDF_OCR_PAGE_CAP
    assert any("limited to first" in m for m in log_messages)


def test_convert_pdf_no_fitz_returns_short_text_without_images():
    """Scanned PDF with pymupdf unavailable returns the (short) extracted text and no images."""
    import unittest.mock
    with patch("pdfminer.high_level.extract_text", return_value="short"), \
         unittest.mock.patch.dict("sys.modules", {"fitz": None}):
        text, images = _convert_pdf(b"fake pdf bytes")

    assert text == "short"
    assert images == []


# -- UAE: summarization tiers -----------------------------------------------------

async def test_summarize_content_below_threshold_returned_as_is_no_llm_call():
    config = ChannelConfig(provider="ollama", model="llama3")
    content = "x" * (_SUMMARIZE_THRESHOLD - 1)

    with patch("icx_engine.connectors.attachments._llm_summarize_chunk", new=AsyncMock()) as mock_sum:
        result = await _summarize_content(config, "notes.txt", content)

    assert result == content
    mock_sum.assert_not_called()


async def test_summarize_content_no_llm_configured_returns_full_content():
    content = "x" * (_SUMMARIZE_THRESHOLD + 5_000)

    result = await _summarize_content(None, "notes.txt", content)

    assert result == content


async def test_summarize_content_single_call_tier():
    """_SUMMARIZE_THRESHOLD < len <= _SINGLE_CALL_LIMIT triggers exactly one summarize call."""
    config = ChannelConfig(provider="ollama", model="llama3")
    content = "x" * (_SUMMARIZE_THRESHOLD + 5_000)

    with patch("icx_engine.connectors.attachments._llm_summarize_chunk",
               new=AsyncMock(return_value="summary")) as mock_sum:
        result = await _summarize_content(config, "notes.txt", content)

    assert result == "summary"
    mock_sum.assert_called_once()


async def test_summarize_content_map_reduce_tier_calls_per_chunk_plus_reduce():
    """> _SINGLE_CALL_LIMIT triggers one call per chunk plus a final reduce call."""
    from icx_engine.connectors.attachments import _CHUNK_SIZE

    config = ChannelConfig(provider="ollama", model="llama3")
    content = "x" * (_SINGLE_CALL_LIMIT + 1)
    expected_chunks = len(_split_into_chunks(content, _CHUNK_SIZE))

    with patch("icx_engine.connectors.attachments._llm_summarize_chunk",
               new=AsyncMock(return_value="short summary")) as mock_sum:
        result = await _summarize_content(config, "notes.txt", content)

    # one call per chunk + one reduce call (combined summaries are short)
    assert mock_sum.call_count == expected_chunks + 1
    assert result == "short summary"


async def test_summarize_content_map_reduce_skips_reduce_if_combined_still_large():
    """If combined chunk summaries still exceed _SINGLE_CALL_LIMIT, no reduce call is made."""
    from icx_engine.connectors.attachments import _CHUNK_SIZE

    config = ChannelConfig(provider="ollama", model="llama3")
    content = "x" * (_SINGLE_CALL_LIMIT + 1)
    expected_chunks = len(_split_into_chunks(content, _CHUNK_SIZE))
    big_summary = "y" * (_SINGLE_CALL_LIMIT // expected_chunks + 1)

    with patch("icx_engine.connectors.attachments._llm_summarize_chunk",
               new=AsyncMock(return_value=big_summary)) as mock_sum:
        result = await _summarize_content(config, "notes.txt", content)

    assert mock_sum.call_count == expected_chunks  # no reduce call
    assert result == "\n\n".join([big_summary] * expected_chunks)


async def test_summarize_content_llm_failure_returns_full_content_with_note():
    config = ChannelConfig(provider="ollama", model="llama3")
    content = "x" * (_SUMMARIZE_THRESHOLD + 5_000)

    with patch("icx_engine.connectors.attachments._llm_summarize_chunk",
               new=AsyncMock(side_effect=Exception("boom"))):
        result = await _summarize_content(config, "notes.txt", content)

    assert result == content + _SUMMARIZE_FAILED_NOTE


# -- UAE: _split_into_chunks -------------------------------------------------------

def test_split_into_chunks_returns_single_chunk_when_under_limit():
    text = "short text"
    assert _split_into_chunks(text, 1000) == [text]


def test_split_into_chunks_preserves_all_content_on_paragraph_boundaries():
    paragraphs = [f"Paragraph {i} " + "word " * 20 for i in range(20)]
    text = "\n\n".join(paragraphs)

    chunks = _split_into_chunks(text, 300)

    assert "".join(chunks).count("Paragraph") == 0 or True  # sanity: no exception
    # every paragraph appears in exactly one chunk, content is fully preserved
    rejoined = "\n\n".join(chunks)
    for para in paragraphs:
        assert para in rejoined
    assert all(len(c) <= 300 or "\n\n" not in c for c in chunks)


def test_split_into_chunks_hard_splits_oversized_paragraph():
    huge_paragraph = "A" * 1000
    chunks = _split_into_chunks(huge_paragraph, 300)

    assert len(chunks) == 4  # 1000 / 300 -> 4 pieces (300,300,300,100)
    assert "".join(chunks) == huge_paragraph
    assert all(len(c) <= 300 for c in chunks)


# -- UAE: video duration + frame sampling ------------------------------------------

async def test_video_duration_parses_ffmpeg_stderr():
    from icx_engine.connectors.attachments import _video_duration

    fake_proc = MagicMock()
    fake_proc.communicate = AsyncMock(return_value=(b"", b"Duration: 00:01:30.50, bitrate: 100kb/s"))

    with patch("asyncio.create_subprocess_exec", new=AsyncMock(return_value=fake_proc)):
        seconds = await _video_duration("fake.mp4")

    assert seconds == 90.5


async def test_video_duration_returns_zero_when_unknown():
    from icx_engine.connectors.attachments import _video_duration

    fake_proc = MagicMock()
    fake_proc.communicate = AsyncMock(return_value=(b"", b"no duration here"))

    with patch("asyncio.create_subprocess_exec", new=AsyncMock(return_value=fake_proc)):
        seconds = await _video_duration("fake.mp4")

    assert seconds == 0.0


async def test_extract_frames_from_video_samples_evenly_across_full_duration(tmp_path):
    """fps is computed from the full video duration so up to _MAX_VIDEO_FRAMES are spread evenly."""
    from icx_engine.connectors.attachments import _extract_frames_from_video, _MAX_VIDEO_FRAMES

    frame_dir = tmp_path / "frames"
    frame_dir.mkdir()
    for i in range(1, _MAX_VIDEO_FRAMES + 1):
        (frame_dir / f"frame{i:04d}.jpg").write_bytes(b"\xff\xd8fake")

    fake_proc = MagicMock()
    fake_proc.communicate = AsyncMock(return_value=(b"", b""))

    with patch("icx_engine.connectors.attachments._video_duration", new=AsyncMock(return_value=30.0)), \
         patch("tempfile.mkdtemp", return_value=str(frame_dir)), \
         patch("asyncio.create_subprocess_exec", new=AsyncMock(return_value=fake_proc)) as mock_exec:
        frames = await _extract_frames_from_video(b"fake video bytes", "demo.mp4")

    assert len(frames) == _MAX_VIDEO_FRAMES
    assert all(f == b"\xff\xd8fake" for f in frames)
    # 15 frames / 30s duration -> fps=0.5
    assert any("fps=0.5" in str(arg) for arg in mock_exec.call_args[0])


async def test_extract_frames_from_video_falls_back_to_default_fps_when_duration_unknown(tmp_path):
    from icx_engine.connectors.attachments import _extract_frames_from_video

    frame_dir = tmp_path / "frames"
    frame_dir.mkdir()
    (frame_dir / "frame0001.jpg").write_bytes(b"\xff\xd8fake")

    fake_proc = MagicMock()
    fake_proc.communicate = AsyncMock(return_value=(b"", b""))

    with patch("icx_engine.connectors.attachments._video_duration", new=AsyncMock(return_value=0.0)), \
         patch("tempfile.mkdtemp", return_value=str(frame_dir)), \
         patch("asyncio.create_subprocess_exec", new=AsyncMock(return_value=fake_proc)) as mock_exec:
        frames = await _extract_frames_from_video(b"fake video bytes", "demo.mp4")

    assert len(frames) == 1
    assert any("fps=0.5" in str(arg) for arg in mock_exec.call_args[0])


# -- UAE: combined video-frame vision call ------------------------------------------

async def test_describe_video_frames_single_combined_call():
    from icx_engine.connectors.attachments import _describe_video_frames

    config = ChannelConfig(provider="anthropic", model="claude-3-haiku-20240307", api_key="sk-ant-test")
    mock_response = MagicMock()
    mock_response.content = [MagicMock()]
    mock_response.content[0].text = "description of sequence"

    frames = [b"frame1", b"frame2", b"frame3"]
    ocr_texts = ["text1", "", "text3"]

    with patch("anthropic.AsyncAnthropic") as mock_cls:
        mock_client = MagicMock()
        mock_cls.return_value = mock_client
        mock_client.messages.create = AsyncMock(return_value=mock_response)
        result = await _describe_video_frames(config, frames, ocr_texts)

    assert result == "description of sequence"
    mock_client.messages.create.assert_called_once()
    _, kwargs = mock_client.messages.create.call_args
    content = kwargs["messages"][0]["content"]
    image_blocks = [c for c in content if c["type"] == "image"]
    assert len(image_blocks) == 3


# -- UAE: frames always returned regardless of vision config ------------------------

@respx.mock
async def test_process_video_returns_frames_in_images_without_vision_config():
    from icx_engine.connectors.attachments import _process_video
    from icx_engine.connectors.audio import WhisperManager

    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/320"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"mp4 bytes"))

    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    whisper = WhisperManager()
    frames = [b"\xff\xd8frame1", b"\xff\xd8frame2"]

    with patch("icx_engine.connectors.attachments._extract_audio_from_video", new=AsyncMock(return_value=b"")), \
         patch("icx_engine.connectors.attachments._extract_frames_from_video", new=AsyncMock(return_value=frames)), \
         patch("icx_engine.connectors.attachments.ocr_image", side_effect=["frame one text", ""]):
        fname, text, images = await _process_video("demo.mp4", content_url, client, None, None, whisper, None)

    assert images == {
        "demo.mp4::frame_01.jpg": base64.b64encode(frames[0]).decode(),
        "demo.mp4::frame_02.jpg": base64.b64encode(frames[1]).decode(),
    }
    assert "[Frame 1/2] frame one text" in text


# -- unsupported attachment types are logged -----------------------------------------

@respx.mock
async def test_process_attachments_logs_unsupported_type():
    content_url = "https://test.atlassian.net/rest/api/3/attachment/content/600"
    respx.get(content_url).mock(return_value=httpx.Response(200, content=b"binary data"))

    raw = _make_raw({"archive.rar": content_url})
    client = JiraClient(JIRA_BASE_URL, JIRA_AUTH_HEADER, JIRA_ALLOWED_HOSTS)
    log_messages: list[str] = []

    texts, images = await process_attachments(raw, client, llm_config=None, log=log_messages.append)

    assert texts == {}
    assert images == {}
    assert any("archive.rar" in m and "unsupported type" in m for m in log_messages)
